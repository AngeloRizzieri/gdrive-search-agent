import io
import os
import re
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
from google.oauth2.credentials import Credentials
from dotenv import load_dotenv

load_dotenv()

SCOPES = ["https://www.googleapis.com/auth/drive.readonly"]
_drive_service = None  # cached only for CLI / local fallback


def _write_from_env_b64(env_var: str, path: str):
    import base64
    val = os.getenv(env_var)
    if val and not os.path.exists(path):
        with open(path, "wb") as f:
            f.write(base64.b64decode(val))


def get_drive_service(credentials=None):
    """
    If credentials (google.oauth2.credentials.Credentials) are provided,
    build a Drive service for that user — used by the web app (multi-user).
    Otherwise fall back to local file-based auth — used by `python main.py`.
    """
    if credentials is not None:
        return build("drive", "v3", credentials=credentials)

    global _drive_service
    if _drive_service:
        return _drive_service

    creds_path = os.getenv("GOOGLE_CREDENTIALS_PATH", "./credentials.json")
    token_path = os.getenv("GOOGLE_TOKEN_PATH", "./token.json")

    _write_from_env_b64("GOOGLE_CREDENTIALS_B64", creds_path)
    _write_from_env_b64("GOOGLE_TOKEN_B64", token_path)

    creds = None
    if os.path.exists(token_path):
        creds = Credentials.from_authorized_user_file(token_path, SCOPES)

    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
        else:
            flow = InstalledAppFlow.from_client_secrets_file(creds_path, SCOPES)
            creds = flow.run_local_server(port=0)
        with open(token_path, "w") as f:
            f.write(creds.to_json())

    _drive_service = build("drive", "v3", credentials=creds)
    return _drive_service


# Only the fields the agent actually needs — keeps tool results lean
_FILE_FIELDS = (
    "id, name, mimeType, modifiedTime, "
    "owners(displayName, emailAddress), "
    "webViewLink"
)


def search_drive(query: str, max_results: int = 10, credentials=None) -> list[dict]:
    service = get_drive_service(credentials)
    all_files = []
    page_token = None
    while len(all_files) < max_results:
        batch = min(1000, max_results - len(all_files))
        kwargs = dict(
            q=f"fullText contains '{query}' and trashed=false",
            pageSize=batch,
            fields=f"nextPageToken,files({_FILE_FIELDS})",
        )
        if page_token:
            kwargs["pageToken"] = page_token
        result = service.files().list(**kwargs).execute()
        all_files.extend(result.get("files", []))
        page_token = result.get("nextPageToken")
        if not page_token:
            break
    return all_files


def list_files(folder_id: str = None, max_results: int = 20, credentials=None) -> list[dict]:
    service = get_drive_service(credentials)
    q = "trashed=false"
    if folder_id:
        q += f" and '{folder_id}' in parents"
    all_files = []
    page_token = None
    while len(all_files) < max_results:
        batch = min(1000, max_results - len(all_files))
        kwargs = dict(
            q=q,
            pageSize=batch,
            fields=f"nextPageToken,files({_FILE_FIELDS})",
        )
        if page_token:
            kwargs["pageToken"] = page_token
        result = service.files().list(**kwargs).execute()
        all_files.extend(result.get("files", []))
        page_token = result.get("nextPageToken")
        if not page_token:
            break
    return all_files


# ── File type handling ────────────────────────────────────────────────────────

# Google Workspace native types → export as text
_EXPORT_MAP = {
    "application/vnd.google-apps.document":     "text/plain",
    "application/vnd.google-apps.spreadsheet":  "text/csv",
    "application/vnd.google-apps.presentation": "text/plain",
    "application/vnd.google-apps.drawing":      "image/svg+xml",
    "application/vnd.google-apps.script":       "application/vnd.google-apps.script+json",
    "application/vnd.google-apps.form":         "application/zip",  # responses export
}

# Text-decodable uploaded types — download raw bytes
_TEXT_TYPES = {
    "text/plain", "text/csv", "text/markdown", "text/html",
    "text/xml", "text/javascript", "text/x-python",
    "application/json", "application/xml",
    "application/javascript", "application/rtf", "text/rtf",
    "image/svg+xml",
}

# Binary types needing a parser
_BINARY_PARSEABLE = {
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",  # .docx
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",        # .xlsx
    "application/vnd.openxmlformats-officedocument.presentationml.presentation", # .pptx
    "application/vnd.ms-powerpoint",                                             # .ppt (fallback)
}

# Truly binary — metadata only
_BINARY_OPAQUE = {
    "image/jpeg", "image/png", "image/gif", "image/webp", "image/bmp", "image/tiff",
    "image/heic", "image/heif",
    "audio/mpeg", "audio/wav", "audio/ogg", "audio/mp4",
    "video/mp4", "video/mpeg", "video/quicktime", "video/x-msvideo",
    "application/zip", "application/x-tar", "application/x-gzip",
    "application/octet-stream",
}

# Image types Claude Vision accepts natively
_OCR_SUPPORTED = {"image/jpeg", "image/png", "image/gif", "image/webp"}

# Image types that need conversion to JPEG before Vision (e.g. iPhone HEIC photos)
_CONVERT_FOR_OCR = {"image/heic", "image/heif", "image/bmp", "image/tiff"}

_MAX_IMAGE_BYTES = 5 * 1024 * 1024  # 5 MB


def _download_bytes(service, file_id: str) -> bytes:
    request = service.files().get_media(fileId=file_id)
    buf = io.BytesIO()
    downloader = MediaIoBaseDownload(buf, request)
    done = False
    while not done:
        _, done = downloader.next_chunk()
    return buf.getvalue()


def _parse_bytes(data: bytes, mime: str) -> str:
    if mime == "application/pdf":
        import pypdf
        reader = pypdf.PdfReader(io.BytesIO(data))
        pages, total = [], 0
        for page in reader.pages:
            t = page.extract_text() or ""
            pages.append(t)
            total += len(t)
            if total >= 3000:  # stop once we have enough for the content limit
                break
        return "\n".join(pages)

    if mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        import docx
        doc = docx.Document(io.BytesIO(data))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    if mime == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        rows = []
        for sheet in wb.worksheets:
            rows.append(f"=== {sheet.title} ===")
            for row in sheet.iter_rows(values_only=True):
                # Skip entirely empty rows (all None or blank strings)
                if not any(v is not None and str(v).strip() for v in row):
                    continue
                rows.append(",".join("" if v is None else str(v) for v in row))
        return "\n".join(rows)

    if mime in ("application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "application/vnd.ms-powerpoint"):
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data))
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
                slides.append(f"[Slide {i}] " + " | ".join(texts))
            return "\n".join(slides)
        except ImportError:
            return "[PPTX parsing unavailable — install python-pptx]"

    return data.decode("utf-8", errors="replace")


def _convert_to_jpeg(data: bytes) -> bytes:
    """Convert any Pillow-supported image (HEIC, BMP, TIFF, …) to JPEG bytes."""
    try:
        import pillow_heif
        pillow_heif.register_heif_opener()
    except ImportError:
        pass  # non-HEIC formats don't need this
    from PIL import Image
    img = Image.open(io.BytesIO(data))
    out = io.BytesIO()
    img.convert("RGB").save(out, format="JPEG", quality=90)
    return out.getvalue()


def _ocr_image(data: bytes, mime: str) -> str:
    """Use Claude Haiku Vision to extract/transcribe text from an image."""
    import base64
    import anthropic
    client = anthropic.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))
    b64 = base64.standard_b64encode(data).decode("utf-8")
    response = client.messages.create(
        model="claude-haiku-4-5-20251001",
        max_tokens=1024,
        messages=[{
            "role": "user",
            "content": [
                {"type": "image", "source": {"type": "base64", "media_type": mime, "data": b64}},
                {"type": "text", "text": (
                    "Extract and transcribe all text visible in this image. "
                    "If it's a handwritten note, transcribe it faithfully. "
                    "If there is little or no text, briefly describe what the image shows."
                )},
            ],
        }],
    )
    return response.content[0].text


def _read_notability(data: bytes) -> str:
    """Read a Notability .note file (ZIP archive containing a PDF)."""
    import zipfile
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as z:
            pdf_names = [n for n in z.namelist() if n.lower().endswith(".pdf")]
            if pdf_names:
                pdf_data = z.read(pdf_names[0])
                return _parse_bytes(pdf_data, "application/pdf")
            txt_names = [n for n in z.namelist() if n.lower().endswith((".txt", ".xml"))]
            if txt_names:
                return z.read(txt_names[0]).decode("utf-8", errors="replace")
        return "[Notability file: no readable content found inside archive]"
    except zipfile.BadZipFile:
        return "[File has .note extension but is not a valid ZIP archive]"
    except Exception as e:
        return f"[Could not read Notability file: {e}]"


def read_document(file_id: str, credentials=None) -> str:
    service = get_drive_service(credentials)
    meta = service.files().get(fileId=file_id, fields=_FILE_FIELDS).execute()
    mime = meta.get("mimeType", "")
    filename = meta.get("name", "")

    owner = meta.get("owners", [{}])[0].get("displayName", "unknown")
    header = (
        f"File: {meta.get('name')}\n"
        f"Modified: {meta.get('modifiedTime', 'unknown')} | Owner: {owner}\n"
        f"---\n"
    )

    # Notability .note files are ZIP archives containing a PDF
    if filename.lower().endswith(".note"):
        try:
            data = _download_bytes(service, file_id)
            text = _read_notability(data)
        except Exception as e:
            text = f"[Error reading Notability file: {e}]"
        return header + text[:6000]

    if mime in _BINARY_OPAQUE:
        if mime in _OCR_SUPPORTED or mime in _CONVERT_FOR_OCR:
            try:
                raw = _download_bytes(service, file_id)
                if len(raw) > _MAX_IMAGE_BYTES:
                    return header + f"[Image too large for analysis ({len(raw) // 1024} KB)]"
                if mime in _CONVERT_FOR_OCR:
                    raw = _convert_to_jpeg(raw)
                    ocr_mime = "image/jpeg"
                else:
                    ocr_mime = mime
                text = _ocr_image(raw, ocr_mime)
                text = re.sub(r'[ \t]+$', '', text, flags=re.MULTILINE)
                text = re.sub(r'\n{3,}', '\n\n', text).strip()
            except Exception as e:
                return header + f"[Image analysis failed: {e}]"
            return header + text[:6000]
        return header + "[Binary file — content not extractable.]"

    try:
        if mime in _EXPORT_MAP:
            response = service.files().export(
                fileId=file_id, mimeType=_EXPORT_MAP[mime]
            ).execute()
            text = response.decode("utf-8") if isinstance(response, bytes) else str(response)
        elif mime in _TEXT_TYPES:
            data = _download_bytes(service, file_id)
            text = data.decode("utf-8", errors="replace")
        elif mime in _BINARY_PARSEABLE:
            data = _download_bytes(service, file_id)
            text = _parse_bytes(data, mime)
        else:
            # Unknown type — attempt raw text decode as last resort
            try:
                data = _download_bytes(service, file_id)
                text = data.decode("utf-8", errors="replace")
            except Exception:
                return header + f"[Unsupported file type: {mime}]"
    except Exception as e:
        return header + f"[Error reading content: {e}]"

    # Collapse whitespace noise: trailing spaces per line, then 3+ blank lines → 1
    text = re.sub(r'[ \t]+$', '', text, flags=re.MULTILINE)
    text = re.sub(r'\n{3,}', '\n\n', text).strip()
    return header + text[:6000]


# ── Formatting ────────────────────────────────────────────────────────────────

def _fmt_file(f: dict) -> str:
    owner = f.get("owners", [{}])[0].get("displayName", "unknown")
    return (
        f"- {f['name']}\n"
        f"  id: {f['id']}\n"
        f"  type: {f.get('mimeType', 'unknown')}\n"
        f"  modified: {f.get('modifiedTime', 'unknown')} | owner: {owner}\n"
        f"  link: {f.get('webViewLink', 'N/A')}"
    )


_MAX_LISTING_CHARS = 4000  # hard cap — prevents a large max_results from flooding context


def execute_tool(name: str, inputs: dict, credentials=None) -> str:
    if name == "search_drive":
        results = search_drive(credentials=credentials, **inputs)
        if not results:
            return "No files found."
        return "\n\n".join(_fmt_file(f) for f in results)[:_MAX_LISTING_CHARS]
    elif name == "list_files":
        results = list_files(credentials=credentials, **inputs)
        if not results:
            return "No files found."
        return "\n\n".join(_fmt_file(f) for f in results)[:_MAX_LISTING_CHARS]
    elif name == "read_document":
        return read_document(credentials=credentials, **inputs)
    else:
        return f"Unknown tool: {name}"


TOOL_SCHEMAS = [
    {
        "name": "search_drive",
        "description": "Search Google Drive by full-text query. Returns id, name, type, modified, owner, link.",
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {"type": "string", "description": "Search query."},
                "max_results": {"type": "integer", "description": "Max results (default 10).", "default": 10},
            },
            "required": ["query"],
        },
    },
    {
        "name": "list_files",
        "description": "List Google Drive files. Optionally filter by folder_id. Returns id, name, type, modified, owner, link.",
        "input_schema": {
            "type": "object",
            "properties": {
                "folder_id": {"type": "string", "description": "Folder ID to list. Omit for all Drive files."},
                "max_results": {"type": "integer", "description": "Max files (default 20).", "default": 20},
            },
            "required": [],
        },
    },
    {
        "name": "read_document",
        "description": "Read a Google Drive file's content by ID. Supports Docs, Sheets, Slides, PDF, DOCX, XLSX, PPTX, plain text, CSV, HTML, JSON, images (OCR), and Notability .note files. Returns up to 6000 chars of content.",
        "input_schema": {
            "type": "object",
            "properties": {
                "file_id": {"type": "string", "description": "Google Drive file ID."},
            },
            "required": ["file_id"],
        },
        # Cache breakpoint on the last tool covers the entire tools array.
        # On every subsequent turn (and question with the same prompt), this
        # segment is billed at ~10% of normal. No effect if below the 1024-token
        # minimum; the API silently ignores it in that case.
        "cache_control": {"type": "ephemeral"},
    },
]
